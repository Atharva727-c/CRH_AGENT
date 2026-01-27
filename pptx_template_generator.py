from __future__ import annotations

import base64
import re
from dataclasses import dataclass
from io import BytesIO
from typing import Iterable
from urllib.parse import quote

import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.xmlchemy import OxmlElement


@dataclass(frozen=True)
class SlideSpec:
    title: str
    content: list[str]
    image_prompt: str


# Default prompts (can be overridden by parsed content)
DEFAULT_SLIDES: list[SlideSpec] = [
    SlideSpec(
        title="1) Executive Summary",
        content=[],
        image_prompt=(
            "corporate executive summary business strategy minimalist design "
            "professional"
        ),
    ),
    SlideSpec(
        title="2) Key Insights",
        content=[],
        image_prompt=(
            "financial growth chart rising bar graph clean background 3d "
            "render"
        ),
    ),
    SlideSpec(
        title="3) Key Trends",
        content=[],
        image_prompt=(
            "global market trends abstract business concept futuristic "
            "blue tone"
        ),
    ),
    SlideSpec(
        title="4) Risks / Watchouts",
        content=[],
        image_prompt=(
            "storm clouds over city skyline business risk concept cinematic "
            "lighting"
        ),
    ),
    SlideSpec(
        title="5) Recommended Next Actions",
        content=[],
        image_prompt=(
            "chess pieces strategy move strategic planning business success"
        ),
    ),
]

CONTENT_FONT_PT = 20
SUBHEADING_FONT_PT = 22


def _strip_markdown_emphasis(text: str) -> str:
    """
    EPAM Dial outputs markdown emphasis like **bold** / *italic*.
    PowerPoint won't render markdown, so we remove the asterisks.
    """
    s = (text or "")
    # Paired emphasis first
    s = re.sub(r"\*\*(.+?)\*\*", r"\1", s)
    s = re.sub(r"\*(.+?)\*", r"\1", s)
    # Remove any remaining stray asterisks
    s = s.replace("**", "").replace("*", "")
    # Normalize whitespace (EPAM sometimes includes double-space line breaks)
    s = re.sub(r"[ \t]+$", "", s)
    return s.strip()


def _split_markdown_bold_heading(line: str) -> tuple[str, str] | None:
    """
    If line begins with markdown bold heading like:
      **Heading**: rest...
    return (Heading, rest). Otherwise None.
    """
    m = re.match(r"^\s*\*\*(.+?)\*\*\s*:?\s*(.*)\s*$", line or "")
    if not m:
        return None
    heading = (m.group(1) or "").strip().rstrip(":")
    rest = (m.group(2) or "").strip()
    if not heading:
        return None
    return heading, rest


def _remove_bullets(paragraph) -> None:
    """
    Remove bullet from a paragraph
    (works for common placeholder bullet styles).
    Uses the underlying oxml since python-pptx doesn't expose a first-class
    API.
    """
    pPr = paragraph._p.get_or_add_pPr()  # noqa: SLF001 (python-pptx internal)
    buNone = OxmlElement("a:buNone")
    pPr.insert(0, buNone)


def _pollinations_url(prompt: str) -> str:
    enc = quote(prompt, safe="")
    return (
        f"https://image.pollinations.ai/prompt/{enc}"
        f"?width=1024&height=1024&nologo=true"
    )


def fetch_pollinations_image(prompt: str, timeout_s: float = 30.0) -> bytes:
    url = _pollinations_url(prompt)
    resp = requests.get(url, timeout=timeout_s)
    resp.raise_for_status()
    return resp.content


def parse_llm_text_to_slides(llm_text: str) -> list[SlideSpec]:
    """
    Parse EPAM Dial narrative output into slide specs.

    Expected headings (but tolerant):
      1) Executive Summary
      2) Key Insights
      3) Key Trends
      4) Risks / Watchouts
      5) Recommended Next Actions

    Bullets are recognized by leading '-' or '•'.
    """
    text = (llm_text or "").strip()
    if not text:
        return DEFAULT_SLIDES

    # Map normalized heading -> SlideSpec template
    templates = {
        re.sub(r"\s+", " ", s.title.lower()).strip(): s for s in DEFAULT_SLIDES
    }

    # Capture numbered headings like "1) Executive Summary" or
    # "### 1) Executive Summary"
    heading_re = re.compile(
        r"^\s*(?:#{1,6}\s*)?(\d+\)\s*.+?)\s*$",
        re.IGNORECASE,
    )
    bullet_re = re.compile(r"^\s*(?:[-•]|\*)\s+(.*)$")

    current_key: str | None = None
    buckets: dict[str, list[str]] = {k: [] for k in templates.keys()}

    for raw_line in text.splitlines():
        line = raw_line.strip()
        if not line:
            continue

        m_head = heading_re.match(line)
        if m_head:
            head = re.sub(r"\s+", " ", m_head.group(1)).strip()
            head = _strip_markdown_emphasis(head)
            norm = head.lower()
            # Find best template match by prefix (e.g. "1) Executive Summary")
            match = None
            for k in templates.keys():
                if norm.startswith(k.split(")")[0] + ")"):
                    match = k
                    break
            if match is None:
                # try direct match
                match = templates.get(norm) and norm
            current_key = match if match in templates else None
            continue

        m_bullet = bullet_re.match(line)
        if m_bullet and current_key:
            buckets[current_key].append(
                _strip_markdown_emphasis(m_bullet.group(1).strip())
            )
            continue

        # Non-bullet lines: treat as continuation for current section
        if current_key:
            buckets[current_key].append(_strip_markdown_emphasis(line))

    out: list[SlideSpec] = []
    for tmpl_key, tmpl in templates.items():
        content = [c for c in buckets.get(tmpl_key, []) if c]
        out.append(
            SlideSpec(
                title=tmpl.title,
                content=content,
                image_prompt=tmpl.image_prompt,
            )
        )
    return out


def generate_pptx_bytes_from_template(
    *,
    template_path: str,
    slides_data: Iterable[SlideSpec],
) -> bytes:
    prs = Presentation(template_path)

    for spec in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title & Content

        # Title
        if slide.shapes.title:
            slide.shapes.title.text = _strip_markdown_emphasis(spec.title)

        # Content placeholder (usually index 1)
        body = None
        if len(slide.placeholders) > 1:
            body = slide.placeholders[1]
        else:
            # fallback: first placeholder that has a text frame and isn't title
            for shp in slide.shapes:
                if (
                    getattr(shp, "has_text_frame", False)
                    and shp != slide.shapes.title
                ):
                    body = shp
                    break

        if body is not None and getattr(body, "has_text_frame", False):
            # Resize to left 60% (approx 5.5 inches width)
            body.left = Inches(0.5)
            body.top = Inches(1.6)
            body.width = Inches(5.5)
            body.height = Inches(5.2)

            tf = body.text_frame
            tf.clear()
            tf.word_wrap = True
            first_paragraph_used = False
            for line in (spec.content or []):
                line = (line or "").strip()
                if not line:
                    continue

                split = _split_markdown_bold_heading(line)
                if split:
                    heading, rest = split
                    heading = _strip_markdown_emphasis(heading)
                    rest = _strip_markdown_emphasis(rest)

                    # Subheading (no bullet)
                    p_head = (
                        tf.paragraphs[0]
                        if not first_paragraph_used
                        else tf.add_paragraph()
                    )
                    first_paragraph_used = True
                    p_head.text = heading
                    p_head.level = 0
                    _remove_bullets(p_head)
                    p_head.font.bold = True
                    p_head.font.size = Pt(SUBHEADING_FONT_PT)
                    for run in p_head.runs:
                        run.font.bold = True
                        run.font.size = Pt(SUBHEADING_FONT_PT)

                    # Optional remainder as a normal bullet under the
                    # subheading.
                    if rest:
                        p_body = tf.add_paragraph()
                        p_body.text = rest
                        p_body.level = 1
                        p_body.font.size = Pt(CONTENT_FONT_PT)
                        for run in p_body.runs:
                            run.font.size = Pt(CONTENT_FONT_PT)
                    continue

                # Normal bullet line
                p = (
                    tf.paragraphs[0]
                    if not first_paragraph_used
                    else tf.add_paragraph()
                )
                first_paragraph_used = True
                p.text = _strip_markdown_emphasis(line)
                p.level = 0
                # Ensure body text isn't too large
                # (template defaults can be ~28pt).
                p.font.size = Pt(CONTENT_FONT_PT)
                for run in p.runs:
                    run.font.size = Pt(CONTENT_FONT_PT)

        # Image on right side
        try:
            img_bytes = fetch_pollinations_image(spec.image_prompt)
            img_stream = BytesIO(img_bytes)
            slide.shapes.add_picture(
                img_stream,
                Inches(6.0),
                Inches(2.0),
                width=Inches(4.0),
            )
        except Exception:
            # If image fails, just skip it (deck still generated)
            pass

    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out.getvalue()


def pptx_bytes_to_data_uri(pptx_bytes: bytes) -> str:
    b64 = base64.b64encode(pptx_bytes).decode("ascii")
    return (
        "data:application/vnd.openxmlformats-officedocument."
        "presentationml.presentation;base64,"
        f"{b64}"
    )
