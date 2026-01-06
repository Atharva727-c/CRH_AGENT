from __future__ import annotations

from functools import lru_cache
from io import BytesIO
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

# Maximum lines allowed per slide to avoid overflow
MAX_LINES_PER_SLIDE = 7


def create_answer_slide(prs: Presentation, title_text: str):
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(245, 247, 250)

    slide.shapes.title.text = title_text
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 58, 138)

    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()

    return tf


def generate_ppt_bytes(title: str, content: str) -> bytes:
    prs = Presentation()

    # ===============================
    # Slide 1: Title Slide
    # ===============================
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.background.fill.solid()
    title_slide.background.fill.fore_color.rgb = RGBColor(30, 58, 138)

    title_slide.shapes.title.text = title
    title_slide.placeholders[1].text = ""

    tp = title_slide.shapes.title.text_frame.paragraphs[0]
    tp.font.size = Pt(36)
    tp.font.bold = True
    tp.font.color.rgb = RGBColor(255, 255, 255)

    # ===============================
    # Slide 2+: Answer Slides
    # ===============================
    tf = create_answer_slide(prs, "Answer")
    line_count = 0

    for raw_line in content.split("\n"):
        line = raw_line.strip()

        # Skip empty lines (prevents empty bullets)
        if not line:
            continue

        # Create new slide if current one is full
        if line_count >= MAX_LINES_PER_SLIDE:
            tf = create_answer_slide(prs, "Answer (cont.)")
            line_count = 0

        # -------- Sources heading --------
        if line.lower().startswith("sources"):
            p = tf.add_paragraph()
            p.text = "Sources"
            p.font.bold = True
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(60, 60, 60)
            p.level = 0
            line_count += 1
            continue

        # -------- Bullet (Streamlit parity) --------
        if line.startswith("- "):
            p = tf.add_paragraph()
            p.text = line[2:].strip()
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.level = 1
            line_count += 1
            continue

        # -------- Normal paragraph --------
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.level = 0
        line_count += 1

    ppt_buffer = BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer.getvalue()


@lru_cache(maxsize=128)
def generate_ppt_bytes_cached(title: str, content: str) -> bytes:
    # Cache by title/content strings to avoid regenerating for the same Q/A.
    return generate_ppt_bytes(title=title, content=content)


def build_ppt_content(answer_text: str | None, sources: Iterable[dict] | None) -> str:
    parts: list[str] = [answer_text.strip() if answer_text else ""]
    src_list = list(sources or [])

    if src_list:
        parts.append("")
        parts.append("Sources:")
        for s in src_list:
            url = (s.get("url") or "").strip()
            title = (s.get("title") or url).strip()
            if not url and not title:
                continue
            # Bullet lines for PPT formatting
            parts.append(f"- {title if title else url}")
            if title and url and url != title:
                parts.append(f"- {url}")

    return "\n".join([p for p in parts if p is not None])


